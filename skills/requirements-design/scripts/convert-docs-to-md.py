#!/usr/bin/env python3
"""
Document to Markdown Converter

Convert various document formats to Markdown for easier access and version control.

Usage:
    # Convert single file
    python convert-docs-to-md.py --file requirements.docx

    # Convert entire folder
    python convert-docs-to-md.py --folder ./docs/Reference/Original/ --output ./docs/Reference/

    # Recursive folder conversion
    python convert-docs-to-md.py --folder ./docs/ --output ./docs/Markdown/ --recursive

Supported formats:
    - .docx (Word documents) → Markdown
    - .pdf (PDF documents) → Markdown (text extraction)
    - .xlsx (Excel spreadsheets) → Markdown tables
    - .pptx (PowerPoint presentations) → Markdown (slide text)
    - .odt (OpenDocument text) → Markdown
    - .rtf (Rich Text Format) → Markdown

Requirements:
    - pandoc (primary converter)
    - Optional: pdftotext (for PDFs), openpyxl (for Excel)
"""

import argparse
import os
import subprocess
import sys
from pathlib import Path
from typing import List, Optional, Tuple


class DocumentConverter:
    """Convert documents to Markdown format"""

    SUPPORTED_FORMATS = {
        '.docx': 'Word Document',
        '.doc': 'Word Document (old format)',
        '.pdf': 'PDF Document',
        '.xlsx': 'Excel Spreadsheet',
        '.xls': 'Excel Spreadsheet (old format)',
        '.pptx': 'PowerPoint Presentation',
        '.ppt': 'PowerPoint Presentation (old format)',
        '.odt': 'OpenDocument Text',
        '.rtf': 'Rich Text Format',
        '.html': 'HTML Document',
        '.htm': 'HTML Document',
    }

    def __init__(self, keep_original: bool = True, overwrite: bool = False,
                 extract_images: bool = False, verbose: bool = True):
        self.keep_original = keep_original
        self.overwrite = overwrite
        self.extract_images = extract_images
        self.verbose = verbose
        self.pandoc_available = self._check_pandoc()

    def _check_pandoc(self) -> bool:
        """Check if pandoc is installed"""
        try:
            result = subprocess.run(['pandoc', '--version'],
                                   capture_output=True, check=False)
            return result.returncode == 0
        except FileNotFoundError:
            return False

    def _print(self, message: str, force: bool = False):
        """Print message if verbose mode enabled"""
        if self.verbose or force:
            print(message)

    def _run_command(self, cmd: List[str]) -> Tuple[bool, str]:
        """Run command and return (success, output)"""
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, check=False)
            if result.returncode == 0:
                return True, result.stdout
            else:
                return False, result.stderr
        except Exception as e:
            return False, str(e)

    def convert_file(self, input_file: Path, output_file: Optional[Path] = None) -> bool:
        """
        Convert single file to Markdown.

        Args:
            input_file: Path to input file
            output_file: Path to output .md file (optional, auto-generated if None)

        Returns:
            True if conversion successful, False otherwise
        """
        if not input_file.exists():
            self._print(f"❌ File not found: {input_file}", force=True)
            return False

        # Check if supported format
        ext = input_file.suffix.lower()
        if ext not in self.SUPPORTED_FORMATS:
            self._print(f"⚠️  Unsupported format: {ext} ({input_file.name})", force=True)
            self._print(f"   Supported: {', '.join(self.SUPPORTED_FORMATS.keys())}", force=True)
            return False

        # Determine output file
        if output_file is None:
            output_file = input_file.with_suffix('.md')
        else:
            # Ensure output has .md extension
            if output_file.suffix != '.md':
                output_file = output_file.with_suffix('.md')

        # Check if output exists and overwrite not enabled
        if output_file.exists() and not self.overwrite:
            self._print(f"⏭️  Skipping (already exists): {output_file.name}")
            return True

        self._print(f"🔄 Converting: {input_file.name} → {output_file.name}")

        # Convert based on file type
        if ext in ['.docx', '.doc', '.odt', '.rtf', '.html', '.htm']:
            return self._convert_with_pandoc(input_file, output_file)
        elif ext == '.pdf':
            return self._convert_pdf(input_file, output_file)
        elif ext in ['.xlsx', '.xls']:
            return self._convert_excel(input_file, output_file)
        elif ext in ['.pptx', '.ppt']:
            return self._convert_powerpoint(input_file, output_file)
        else:
            self._print(f"❌ No converter available for: {ext}", force=True)
            return False

    def _convert_with_pandoc(self, input_file: Path, output_file: Path) -> bool:
        """Convert using pandoc"""
        if not self.pandoc_available:
            self._print("❌ Pandoc not installed. Install with:", force=True)
            self._print("   Linux:   sudo apt install pandoc", force=True)
            self._print("   macOS:   brew install pandoc", force=True)
            self._print("   Windows: https://pandoc.org/installing.html", force=True)
            return False

        # Build pandoc command
        cmd = ['pandoc', str(input_file), '-t', 'markdown', '-o', str(output_file)]

        # Add image extraction if enabled
        if self.extract_images:
            images_dir = output_file.parent / 'images'
            images_dir.mkdir(exist_ok=True)
            cmd.extend(['--extract-media', str(images_dir)])

        # Additional pandoc options for better markdown
        cmd.extend([
            '--wrap=none',  # Don't wrap long lines
            '--markdown-headings=atx',  # Use # style headings
        ])

        success, output = self._run_command(cmd)

        if success:
            self._print(f"✅ Converted: {output_file.name}")
            return True
        else:
            self._print(f"❌ Conversion failed: {output}", force=True)
            return False

    def _convert_pdf(self, input_file: Path, output_file: Path) -> bool:
        """Convert PDF to Markdown"""
        # Try pandoc first
        if self.pandoc_available:
            self._print("   Using pandoc for PDF conversion...")
            return self._convert_with_pandoc(input_file, output_file)

        # Fallback to pdftotext
        self._print("   Pandoc not available, trying pdftotext...", force=True)

        # Check if pdftotext available
        success, _ = self._run_command(['which', 'pdftotext'])
        if not success:
            self._print("❌ Neither pandoc nor pdftotext available", force=True)
            self._print("   Install: sudo apt install poppler-utils", force=True)
            return False

        # Extract text with pdftotext
        text_file = input_file.with_suffix('.txt')
        success, output = self._run_command(['pdftotext', str(input_file), str(text_file)])

        if not success:
            self._print(f"❌ PDF extraction failed: {output}", force=True)
            return False

        # Read text and create markdown
        try:
            with open(text_file, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()

            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"# {input_file.stem}\n\n")
                f.write("*Converted from PDF - formatting may be limited*\n\n")
                f.write("---\n\n")
                f.write(text)

            # Remove temporary text file
            text_file.unlink()

            self._print(f"✅ Converted: {output_file.name}")
            return True

        except Exception as e:
            self._print(f"❌ PDF conversion failed: {e}", force=True)
            return False

    def _convert_excel(self, input_file: Path, output_file: Path) -> bool:
        """Convert Excel to Markdown tables"""
        try:
            import openpyxl
        except ImportError:
            self._print("❌ openpyxl not installed. Install with:", force=True)
            self._print("   pip install openpyxl", force=True)
            return False

        try:
            # Load workbook
            wb = openpyxl.load_workbook(input_file, data_only=True)

            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"# {input_file.stem}\n\n")
                f.write("*Converted from Excel*\n\n")
                f.write("---\n\n")

                # Process each sheet
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    f.write(f"## {sheet_name}\n\n")

                    # Get data
                    data = list(sheet.values)
                    if not data:
                        f.write("*Empty sheet*\n\n")
                        continue

                    # Write as markdown table
                    # Header row
                    header = data[0]
                    f.write("| " + " | ".join(str(cell or "") for cell in header) + " |\n")
                    f.write("|" + " --- |" * len(header) + "\n")

                    # Data rows
                    for row in data[1:]:
                        if any(cell for cell in row):  # Skip empty rows
                            f.write("| " + " | ".join(str(cell or "") for cell in row) + " |\n")

                    f.write("\n")

            self._print(f"✅ Converted: {output_file.name}")
            return True

        except Exception as e:
            self._print(f"❌ Excel conversion failed: {e}", force=True)
            return False

    def _convert_powerpoint(self, input_file: Path, output_file: Path) -> bool:
        """Convert PowerPoint to Markdown"""
        # Try pandoc first
        if self.pandoc_available:
            return self._convert_with_pandoc(input_file, output_file)

        # Fallback: try python-pptx
        try:
            from pptx import Presentation
        except ImportError:
            self._print("❌ Neither pandoc nor python-pptx available", force=True)
            self._print("   Install: pip install python-pptx", force=True)
            return False

        try:
            prs = Presentation(input_file)

            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"# {input_file.stem}\n\n")
                f.write("*Converted from PowerPoint*\n\n")
                f.write("---\n\n")

                # Process each slide
                for i, slide in enumerate(prs.slides, 1):
                    f.write(f"## Slide {i}\n\n")

                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            f.write(f"{shape.text}\n\n")

                    f.write("---\n\n")

            self._print(f"✅ Converted: {output_file.name}")
            return True

        except Exception as e:
            self._print(f"❌ PowerPoint conversion failed: {e}", force=True)
            return False

    def convert_folder(self, input_folder: Path, output_folder: Path,
                      recursive: bool = False) -> Tuple[int, int]:
        """
        Convert all supported documents in folder to Markdown.

        Args:
            input_folder: Input folder path
            output_folder: Output folder path
            recursive: Recursively process subfolders

        Returns:
            Tuple of (successful_conversions, failed_conversions)
        """
        if not input_folder.exists():
            self._print(f"❌ Folder not found: {input_folder}", force=True)
            return 0, 0

        # Create output folder
        output_folder.mkdir(parents=True, exist_ok=True)

        success_count = 0
        fail_count = 0

        # Get files to convert
        if recursive:
            files = [f for f in input_folder.rglob('*') if f.is_file()]
        else:
            files = [f for f in input_folder.iterdir() if f.is_file()]

        # Filter supported formats
        supported_files = [f for f in files if f.suffix.lower() in self.SUPPORTED_FORMATS]

        if not supported_files:
            self._print(f"⚠️  No supported documents found in {input_folder}", force=True)
            return 0, 0

        self._print(f"\n📁 Converting {len(supported_files)} files from {input_folder}")
        self._print(f"📁 Output to: {output_folder}\n")

        for file in supported_files:
            # Calculate relative path for subdirectories
            if recursive:
                rel_path = file.relative_to(input_folder)
                output_file = output_folder / rel_path.with_suffix('.md')
                output_file.parent.mkdir(parents=True, exist_ok=True)
            else:
                output_file = output_folder / file.with_suffix('.md').name

            if self.convert_file(file, output_file):
                success_count += 1
            else:
                fail_count += 1

        return success_count, fail_count


def main():
    parser = argparse.ArgumentParser(
        description='Convert documents to Markdown format',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )

    # Info options (check before requiring input)
    parser.add_argument('--supported-formats', action='store_true',
                       help='Show supported formats and exit')

    # Input options (mutually exclusive)
    input_group = parser.add_mutually_exclusive_group(required=False)
    input_group.add_argument('--file', type=Path, help='Convert single file')
    input_group.add_argument('--folder', type=Path, help='Convert all files in folder')

    # Output options
    parser.add_argument('--output', type=Path,
                       help='Output file/folder (default: same location as input)')

    # Behavior options
    parser.add_argument('--recursive', action='store_true',
                       help='Recursively process subfolders (with --folder)')
    parser.add_argument('--overwrite', action='store_true',
                       help='Overwrite existing .md files (default: skip)')
    parser.add_argument('--extract-images', action='store_true',
                       help='Extract images to separate folder (requires pandoc)')
    parser.add_argument('--delete-original', action='store_true',
                       help='Delete original file after successful conversion')
    parser.add_argument('--quiet', action='store_true',
                       help='Suppress output (only errors)')

    args = parser.parse_args()

    # Show supported formats
    if args.supported_formats:
        print("Supported document formats:")
        for ext, desc in DocumentConverter.SUPPORTED_FORMATS.items():
            print(f"  {ext:8} - {desc}")
        print("\nRequired tools:")
        print("  - pandoc (most formats)")
        print("  - pdftotext (PDFs without pandoc)")
        print("  - openpyxl (Excel files) - pip install openpyxl")
        print("  - python-pptx (PowerPoint without pandoc) - pip install python-pptx")
        sys.exit(0)

    # Require --file or --folder if not showing formats
    if not args.file and not args.folder:
        parser.error("one of the arguments --file or --folder is required")

    # Initialize converter
    converter = DocumentConverter(
        keep_original=not args.delete_original,
        overwrite=args.overwrite,
        extract_images=args.extract_images,
        verbose=not args.quiet
    )

    # Convert file or folder
    if args.file:
        # Single file conversion
        success = converter.convert_file(args.file, args.output)
        sys.exit(0 if success else 1)
    else:
        # Folder conversion
        output_folder = args.output or args.folder / 'Markdown'
        success_count, fail_count = converter.convert_folder(
            args.folder, output_folder, args.recursive
        )

        print(f"\n✅ Successfully converted: {success_count}")
        if fail_count > 0:
            print(f"❌ Failed: {fail_count}")

        sys.exit(0 if fail_count == 0 else 1)


if __name__ == '__main__':
    main()
